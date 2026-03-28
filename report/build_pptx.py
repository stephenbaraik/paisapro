"""
PaisaPro.ai — Beautiful PowerPoint Presentation Builder
Generates a polished, visually appealing .pptx for faculty pitch.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
import os

# ── Brand Colors ─────────────────────────────────────────────────────────────
PRIMARY   = RGBColor(15, 23, 42)      # Dark navy
ACCENT    = RGBColor(99, 102, 241)    # Indigo
GREEN     = RGBColor(16, 185, 129)    # Emerald
RED       = RGBColor(239, 68, 68)     # Red
AMBER     = RGBColor(245, 158, 11)    # Amber
CYAN      = RGBColor(6, 182, 212)     # Cyan
WHITE     = RGBColor(255, 255, 255)
LIGHT_BG  = RGBColor(241, 245, 249)   # Slate-100
GRAY      = RGBColor(100, 116, 139)   # Slate-500
DARK_GRAY = RGBColor(51, 65, 85)      # Slate-700
LIGHT_ACCENT = RGBColor(224, 225, 254)  # Indigo-100
LIGHT_GREEN  = RGBColor(209, 250, 229)  # Emerald-100
LIGHT_CYAN   = RGBColor(207, 250, 254)  # Cyan-100
LIGHT_AMBER  = RGBColor(254, 243, 199)  # Amber-100
LIGHT_RED    = RGBColor(254, 226, 226)  # Red-100

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


# ═══════════════════════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, w, h, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def _add_text_box(slide, left, top, w, h, text, font_size=14,
                  color=PRIMARY, bold=False, align=PP_ALIGN.LEFT,
                  font_name="Segoe UI"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tb


def _add_rich_box(slide, left, top, w, h):
    """Return a text_frame ready for multiple paragraphs."""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def _add_para(tf, text, size=13, color=PRIMARY, bold=False, space_after=Pt(4),
              space_before=Pt(0), align=PP_ALIGN.LEFT, font_name="Segoe UI",
              bullet=False):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_after = space_after
    p.space_before = space_before
    if bullet:
        p.level = 0
    return p


def _add_bullet(tf, text, size=12, color=DARK_GRAY, indent_level=0,
                bold=False, space_after=Pt(3)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.level = indent_level
    p.space_after = space_after
    return p


def _section_header(slide, text, subtitle=""):
    """Add consistent slide title bar."""
    # Accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0),
                                 SLIDE_W, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    _add_text_box(slide, Inches(0.7), Inches(0.25), Inches(11), Inches(0.6),
                  text, font_size=28, color=PRIMARY, bold=True)

    if subtitle:
        _add_text_box(slide, Inches(0.7), Inches(0.80), Inches(11), Inches(0.4),
                      subtitle, font_size=14, color=GRAY)


def _card(slide, left, top, w, h, title, body_lines, accent_color,
          title_size=14, body_size=11):
    """A rounded card with colored top border."""
    # Background card
    card = _add_rect(slide, left, top, w, h, WHITE, border_color=RGBColor(226, 232, 240), radius=0.05)
    # Top accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     left + Emu(5000), top + Emu(5000),
                                     w - Emu(10000), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent_color
    stripe.line.fill.background()

    # Title
    _add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                  w - Inches(0.4), Inches(0.35),
                  title, font_size=title_size, color=accent_color, bold=True)

    # Body
    tf = _add_rich_box(slide, left + Inches(0.2), top + Inches(0.5),
                       w - Inches(0.4), h - Inches(0.6))
    for line in body_lines:
        _add_para(tf, line, size=body_size, color=DARK_GRAY, space_after=Pt(3))

    return card


def _kpi_card(slide, left, top, w, h, number, label, accent_color):
    """Big number KPI card."""
    card = _add_rect(slide, left, top, w, h, WHITE, border_color=accent_color, radius=0.08)
    _add_text_box(slide, left, top + Inches(0.2), w, Inches(0.55),
                  number, font_size=36, color=accent_color, bold=True,
                  align=PP_ALIGN.CENTER)
    _add_text_box(slide, left, top + Inches(0.75), w, Inches(0.4),
                  label, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)


def _add_slide():
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, WHITE)
    return slide


def _add_slide_number(slide, num, total):
    _add_text_box(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.35),
                  f"{num} / {total}", font_size=9, color=GRAY,
                  align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════

def slide_title():
    slide = _add_slide()
    _set_bg(slide, PRIMARY)

    # Decorative circles (subtle)
    for cx, cy, sz, op in [(0.5, 0.5, 3, 0.04), (11, 6, 4, 0.03), (10, 1, 2.5, 0.03)]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(cx), Inches(cy),
                                    Inches(sz), Inches(sz))
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT
        c.fill.fore_color.brightness = 0.7
        c.line.fill.background()

    # NMIMS logo
    logo_path = os.path.join(os.path.dirname(__file__), "NMIMS logo.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(5.6), Inches(0.4),
                                  height=Inches(0.9))

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(4.7), Inches(1.55),
                                   Inches(4), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    _add_text_box(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(0.9),
                  "PaisaPro.ai", font_size=52, color=WHITE, bold=True,
                  align=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(2.7), Inches(11.3), Inches(0.5),
                  "AI-Powered Investment Advisory Platform",
                  font_size=22, color=CYAN, align=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(0.4),
                  "for the Indian Equity Market",
                  font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # Subtitle
    _add_text_box(slide, Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.6),
                  "A Full-Stack System Integrating Quantitative Finance, Machine Learning,\nand Large Language Models for Retail Investor Decision Support",
                  font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

    # Team
    _add_text_box(slide, Inches(1), Inches(4.9), Inches(11.3), Inches(0.4),
                  "Stephen Baraik (86032300068)    |    Tanisha Ghosh (86032300074)    |    Sneha Das (86032300029)",
                  font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Supervisor
    _add_text_box(slide, Inches(1), Inches(5.6), Inches(11.3), Inches(0.35),
                  "Under the supervision of Dr. Suresh Pathare",
                  font_size=11, color=GRAY, align=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1), Inches(5.9), Inches(11.3), Inches(0.25),
                  "Associate Professor, Data Science",
                  font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # Footer
    _add_text_box(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.3),
                  "SVKM's NMIMS University  ·  SOMASA  ·  2025–2026",
                  font_size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════════════════════

def slide_agenda():
    slide = _add_slide()
    _section_header(slide, "Agenda")

    items = [
        ("01", "Problem Statement", ACCENT),
        ("02", "Our Solution", ACCENT),
        ("03", "System Architecture", GREEN),
        ("04", "Data Pipeline & Database", GREEN),
        ("05", "Analytics Engine", CYAN),
        ("06", "Machine Learning Pipeline", CYAN),
        ("07", "Time-Series & Volatility", AMBER),
        ("08", "Risk & Portfolio", AMBER),
        ("09", "Agentic AI Advisor", ACCENT),
        ("10", "Tech Stack & Frontend", GREEN),
        ("11", "Results & Validation", CYAN),
        ("12", "Comparison", CYAN),
        ("13", "Limitations & Future Work", AMBER),
        ("14", "Conclusion", ACCENT),
    ]

    cols = 2
    per_col = 7
    for idx, (num, label, color) in enumerate(items):
        col = idx // per_col
        row = idx % per_col
        lx = Inches(1.5) + Inches(col * 5.5)
        ty = Inches(1.5) + Inches(row * 0.72)

        # Number circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, lx, ty,
                                       Inches(0.45), Inches(0.45))
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        tf = circ.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = "Segoe UI"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        _add_text_box(slide, lx + Inches(0.6), ty + Inches(0.05),
                      Inches(4.5), Inches(0.35),
                      label, font_size=15, color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════

def slide_problem():
    slide = _add_slide()
    _section_header(slide, "The Problem",
                    "Why we built PaisaPro.ai")

    # KPI cards
    _kpi_card(slide, Inches(0.8), Inches(1.5), Inches(2.8), Inches(1.3),
              "150M+", "Demat Accounts in India", ACCENT)
    _kpi_card(slide, Inches(4.0), Inches(1.5), Inches(2.8), Inches(1.3),
              "₹20K Cr", "Monthly SIP Inflows", GREEN)
    _kpi_card(slide, Inches(7.2), Inches(1.5), Inches(2.8), Inches(1.3),
              "4+", "Disconnected Tools", RED)

    # Three gaps
    gap_top = Inches(3.2)
    for i, (title, desc, color) in enumerate([
        ("Analytics Gap",
         "ARIMA, GARCH, Fama-French, Markowitz require deep expertise in statistics & programming. No single Indian platform integrates all of these.",
         ACCENT),
        ("AI Advisory Gap",
         "LLMs like ChatGPT give generic advice without live market data, real-time signals, or user portfolio context. Advice is outdated by design.",
         CYAN),
        ("Tool Fragmentation",
         "Retail investors juggle Screener.in (fundamentals) + TradingView (charts) + Groww (trading) + ChatGPT (advice) — 4 disconnected tools.",
         AMBER),
    ]):
        lx = Inches(0.8 + i * 4.1)
        _card(slide, lx, gap_top, Inches(3.7), Inches(2.5),
              title, [desc], color, body_size=11)

    # Bottom callout
    _add_rect(slide, Inches(2), Inches(6.1), Inches(9.3), Inches(0.6),
              LIGHT_RED, border_color=RED, radius=0.1)
    _add_text_box(slide, Inches(2), Inches(6.15), Inches(9.3), Inches(0.5),
                  "No single platform unifies quantitative analytics, ML, and AI advisory for Indian retail investors",
                  font_size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — SOLUTION (FOUR PILLARS)
# ═══════════════════════════════════════════════════════════════════════════════

def slide_solution():
    slide = _add_slide()
    _section_header(slide, "Our Solution: Four Pillars")

    pillars = [
        ("Quantitative\nFinance", ACCENT, LIGHT_ACCENT,
         ["ARIMA forecasting", "GARCH(1,1) volatility", "Fama-French 4-factor", "Markowitz MVO", "Monte Carlo (1K sims)"]),
        ("Machine\nLearning", GREEN, LIGHT_GREEN,
         ["Random Forest classifier", "3-model regression", "Isolation Forest", "K-Means clustering", "Walk-forward CV"]),
        ("Agentic\nAI Advisor", CYAN, LIGHT_CYAN,
         ["Llama 3.3 70B via Groq", "4 live function-calling tools", "Portfolio context injection", "SSE streaming", "Multi-turn reasoning"]),
        ("Full-Stack\nPlatform", AMBER, LIGHT_AMBER,
         ["React 19 + TypeScript", "FastAPI backend", "18 pages · 35+ APIs", "Dark/light themes", "Cloud deployed"]),
    ]

    for i, (title, color, bg_color, items) in enumerate(pillars):
        lx = Inches(0.5 + i * 3.15)
        ty = Inches(1.4)
        w = Inches(2.9)
        h = Inches(4.2)

        card = _add_rect(slide, lx, ty, w, h, bg_color, border_color=color, radius=0.06)

        # Icon circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       lx + Inches(0.95), ty + Inches(0.2),
                                       Inches(0.7), Inches(0.7))
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        icons = ["Σ", "Ω", "AI", "< >"]
        tf = circ.text_frame
        tf.paragraphs[0].text = icons[i]
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        _add_text_box(slide, lx + Inches(0.1), ty + Inches(1.0), w - Inches(0.2), Inches(0.6),
                      title, font_size=16, color=color, bold=True, align=PP_ALIGN.CENTER)

        rtf = _add_rich_box(slide, lx + Inches(0.25), ty + Inches(1.7),
                            w - Inches(0.5), Inches(2.3))
        for item in items:
            _add_para(rtf, f"▸  {item}", size=11, color=DARK_GRAY, space_after=Pt(4))

    # Bottom bar
    _add_rect(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.55),
              PRIMARY, radius=0.08)
    _add_text_box(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.45),
                  "500+ NSE Stocks  ·  Daily Auto-Refresh  ·  Sub-100ms Cached Responses  ·  19 Service Modules",
                  font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════

def slide_architecture():
    slide = _add_slide()
    _section_header(slide, "System Architecture",
                    "Layered monolithic design with 19 independent service modules")

    layers = [
        ("FRONTEND", Inches(1.4), ACCENT, LIGHT_ACCENT,
         "React 19  +  TypeScript  +  Tailwind CSS  +  Recharts  ·  18 Interactive Pages"),
        ("API LAYER", Inches(2.5), GREEN, LIGHT_GREEN,
         "FastAPI  ·  6 Routers  ·  35+ Endpoints  ·  Pydantic v2 Validation  ·  CORS + SSE"),
        ("SERVICES", Inches(3.6), CYAN, LIGHT_CYAN,
         "Technical Analysis  |  ML Models (RF, GBM, EN)  |  Forecasting (ARIMA, GARCH)  |  AI Advisor (Groq)  |  Risk Factors  |  Sector Rotation  |  News Sentiment"),
        ("DATA", Inches(4.7), AMBER, LIGHT_AMBER,
         "Supabase PostgreSQL  ·  yfinance API  ·  Groq LLM API  ·  Google News RSS  ·  In-Memory TTL Cache (5 tiers)"),
    ]

    for label, ty, color, bg, desc in layers:
        # Label
        _add_rect(slide, Inches(0.7), ty, Inches(1.5), Inches(0.7), color, radius=0.06)
        _add_text_box(slide, Inches(0.7), ty + Inches(0.1), Inches(1.5), Inches(0.5),
                      label, font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # Content bar
        _add_rect(slide, Inches(2.4), ty, Inches(10.2), Inches(0.7), bg, border_color=color, radius=0.04)
        _add_text_box(slide, Inches(2.6), ty + Inches(0.1), Inches(9.8), Inches(0.5),
                      desc, font_size=12, color=DARK_GRAY)

    # Arrows between layers
    for y in [Inches(2.1), Inches(3.2), Inches(4.3)]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                        Inches(6.5), y, Inches(0.4), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()

    # Deployment info
    _card(slide, Inches(0.7), Inches(5.8), Inches(3.5), Inches(1.2),
          "Deployment", ["Frontend: Netlify CDN", "Backend: HF Spaces (Docker)", "Database: Supabase Cloud", "LLM: Groq Cloud"],
          GREEN, body_size=11)

    _card(slide, Inches(4.6), Inches(5.8), Inches(3.8), Inches(1.2),
          "Cache Architecture", ["Indices: 5 min TTL", "News/Macro: 30 min", "Analytics: 1–2 hours", "ML Models/Prices: 24 hours"],
          AMBER, body_size=11)

    _card(slide, Inches(8.8), Inches(5.8), Inches(3.8), Inches(1.2),
          "Security", ["CORS whitelist", "Env-var secrets (no client keys)", "Supabase RLS", "Pydantic input validation"],
          RED, body_size=11)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — DATA PIPELINE
# ═══════════════════════════════════════════════════════════════════════════════

def slide_data_pipeline():
    slide = _add_slide()
    _section_header(slide, "Data Pipeline & Database",
                    "Automated daily refresh for 500+ NSE stocks")

    # Left: pipeline steps
    _card(slide, Inches(0.7), Inches(1.4), Inches(5.5), Inches(4.5),
          "Daily Refresh Pipeline (6 PM IST)", [], ACCENT)

    steps = [
        ("1", "Clear all in-memory caches", ACCENT),
        ("2", "Batch-fetch OHLCV via yfinance", GREEN),
        ("3", "Upsert 500+ stocks into Supabase", GREEN),
        ("4", "Compute 30+ technical indicators (SMA, EMA, MACD, RSI, BB, ATR)", CYAN),
        ("5", "Refresh macro data (VIX, USD/INR, Gold, Crude Oil)", CYAN),
        ("6", "Pre-warm analytics report cache", AMBER),
        ("7", "Retrain 500+ Random Forest classifiers", AMBER),
    ]
    for i, (num, text, color) in enumerate(steps):
        ty = Inches(2.0 + i * 0.52)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(1.0), ty, Inches(0.35), Inches(0.35))
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        tf = circ.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        _add_text_box(slide, Inches(1.5), ty, Inches(4.5), Inches(0.35),
                      text, font_size=12, color=DARK_GRAY)

    # Right: database tables
    _card(slide, Inches(6.6), Inches(1.4), Inches(6.0), Inches(2.3),
          "Database Schema (Supabase PostgreSQL)", [
              "stocks — symbol (PK), company_name, sector, current_price",
              "stock_prices — (symbol, date) PK, OHLCV + 30 indicators",
              "macro_prices — (ticker, date) PK, close price",
              "portfolio_holdings — user holdings with buy price/date",
              "watchlist — user watchlist items",
          ], GREEN, body_size=11)

    _card(slide, Inches(6.6), Inches(4.0), Inches(6.0), Inches(1.9),
          "Indexing Strategy", [
              "(symbol, date) UNIQUE — fast lookup + upsert dedup",
              "(symbol) — per-stock history queries",
              "(date) — cross-stock date range queries",
              "(sector) — sector-based filtering",
              "(user_id) — user-specific portfolio queries",
          ], CYAN, body_size=11)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — ANALYTICS ENGINE
# ═══════════════════════════════════════════════════════════════════════════════

def slide_analytics():
    slide = _add_slide()
    _section_header(slide, "Analytics Engine",
                    "7-indicator weighted composite signal generation")

    # Indicators list (left)
    indicators = [
        ("RSI (14)", "25%", "< 30 → BUY (+2),  > 70 → SELL (−2)"),
        ("MACD (12,26,9)", "20%", "Above signal → BUY,  Below → SELL"),
        ("Bollinger Bands (20,2)", "20%", "Below lower → BUY,  Above upper → SELL"),
        ("ATR (14)", "15%", "Low relative → favorable,  High → caution"),
        ("ADX (14)", "10%", "> 25 trending,  +DI vs −DI direction"),
        ("VWAP", "5%", "Price > VWAP → bullish bias"),
        ("OBV", "5%", "Rising OBV → buying pressure confirmation"),
    ]

    for i, (name, weight, logic) in enumerate(indicators):
        ty = Inches(1.5 + i * 0.7)

        # Weight badge
        badge = _add_rect(slide, Inches(0.7), ty, Inches(0.8), Inches(0.45),
                          ACCENT, radius=0.15)
        _add_text_box(slide, Inches(0.7), ty + Inches(0.05), Inches(0.8), Inches(0.35),
                      weight, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        _add_text_box(slide, Inches(1.65), ty, Inches(2.2), Inches(0.45),
                      name, font_size=13, color=PRIMARY, bold=True)
        _add_text_box(slide, Inches(3.9), ty + Inches(0.02), Inches(4), Inches(0.4),
                      logic, font_size=10, color=GRAY)

    # Right side — signal output
    _add_rect(slide, Inches(8.3), Inches(1.4), Inches(4.3), Inches(2.8),
              LIGHT_GREEN, border_color=GREEN, radius=0.06)
    _add_text_box(slide, Inches(8.5), Inches(1.6), Inches(3.9), Inches(0.4),
                  "Composite Signal Output", font_size=16, color=GREEN, bold=True,
                  align=PP_ALIGN.CENTER)
    tf = _add_rich_box(slide, Inches(8.5), Inches(2.1), Inches(3.9), Inches(1.8))
    _add_para(tf, "Score = Σ wᵢ · sᵢ", size=14, color=PRIMARY, bold=True,
              align=PP_ALIGN.CENTER, space_after=Pt(12))
    _add_para(tf, "Score > +0.3  →  BUY", size=13, color=GREEN, bold=True, space_after=Pt(6))
    _add_para(tf, "Score ∈ [−0.3, +0.3]  →  HOLD", size=13, color=AMBER, bold=True, space_after=Pt(6))
    _add_para(tf, "Score < −0.3  →  SELL", size=13, color=RED, bold=True)

    # Distribution box
    _add_rect(slide, Inches(8.3), Inches(4.5), Inches(4.3), Inches(1.8),
              LIGHT_BG, border_color=ACCENT, radius=0.06)
    _add_text_box(slide, Inches(8.5), Inches(4.6), Inches(3.9), Inches(0.35),
                  "Typical Nifty 500 Distribution", font_size=13, color=ACCENT, bold=True,
                  align=PP_ALIGN.CENTER)
    tf2 = _add_rich_box(slide, Inches(8.6), Inches(5.0), Inches(3.8), Inches(1.2))
    _add_para(tf2, "BUY      120–180 stocks    (24–36%)", size=12, color=GREEN, bold=True, space_after=Pt(4))
    _add_para(tf2, "HOLD    200–280 stocks    (40–56%)", size=12, color=AMBER, bold=True, space_after=Pt(4))
    _add_para(tf2, "SELL     80–140 stocks      (16–28%)", size=12, color=RED, bold=True)

    _add_text_box(slide, Inches(0.7), Inches(6.5), Inches(7), Inches(0.35),
                  "Confidence: percentage of indicators agreeing with composite direction (30%–90%+)",
                  font_size=11, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — ML PIPELINE
# ═══════════════════════════════════════════════════════════════════════════════

def slide_ml():
    slide = _add_slide()
    _section_header(slide, "Machine Learning Pipeline",
                    "Classification, regression ensemble, anomaly detection & clustering")

    # Left column — Classification
    _card(slide, Inches(0.5), Inches(1.4), Inches(4.0), Inches(2.8),
          "Random Forest Classifier", [
              "Task: 5-day directional prediction",
              "100 decision trees",
              "252-day rolling training window",
              "Features: RSI, MACD, BB width, vol ratio, price vs SMA20, daily return",
              "",
              "Accuracy: 52–58% (above 50% random baseline)",
          ], ACCENT)

    # Left column — Unsupervised
    _card(slide, Inches(0.5), Inches(4.5), Inches(4.0), Inches(2.4),
          "Unsupervised Models", [
              "K-Means (k=4): Defensive · Value · Momentum · High-Vol",
              "Features: volatility, momentum, beta",
              "",
              "Isolation Forest: anomaly detection",
              "Detects earnings spikes, index rebalancing, panics",
          ], GREEN)

    # Right column — Regression Ensemble
    _card(slide, Inches(4.8), Inches(1.4), Inches(7.8), Inches(5.5),
          "3-Model Regression Ensemble — 30-Day Return Forecast", [
              "",
          ], CYAN)

    # Three model sub-cards
    models = [
        ("Random Forest", "200 trees, depth 6, √p features", ACCENT),
        ("ElasticNet", "L1+L2, RobustScaler, winsorized", GREEN),
        ("HistGradientBoosting", "100 est, LR 0.05, early stopping", AMBER),
    ]
    for i, (name, desc, color) in enumerate(models):
        lx = Inches(5.1 + i * 2.5)
        _add_rect(slide, lx, Inches(2.1), Inches(2.3), Inches(0.9),
                  WHITE, border_color=color, radius=0.06)
        _add_text_box(slide, lx + Inches(0.1), Inches(2.15), Inches(2.1), Inches(0.35),
                      name, font_size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
        _add_text_box(slide, lx + Inches(0.1), Inches(2.5), Inches(2.1), Inches(0.4),
                      desc, font_size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # Key innovations
    tf = _add_rich_box(slide, Inches(5.1), Inches(3.3), Inches(7.2), Inches(3.2))
    _add_para(tf, "Key Innovations:", size=13, color=PRIMARY, bold=True, space_after=Pt(8))
    for item in [
        "32+ features: OHLC microstructure + technicals + extended",
        "Vol-adjusted log-return target: y = ln(Ct+h/Ct) / σ₂₀",
        "Walk-forward CV with embargo gap (prevents look-ahead bias)",
        "Sharpe-weighted ensemble: ŷ = Σ(SRₘ · ŷₘ) / Σ SRₘ",
        "80% confidence interval from RF tree distribution",
        "",
        "Directional accuracy: 55–60% (ensemble > individual models)",
    ]:
        _add_para(tf, f"▸  {item}" if item else "", size=11, color=DARK_GRAY,
                  space_after=Pt(4))


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — FORECASTING
# ═══════════════════════════════════════════════════════════════════════════════

def slide_forecasting():
    slide = _add_slide()
    _section_header(slide, "Time-Series & Volatility Forecasting",
                    "ARIMA, Holt-Winters ETS, and GARCH(1,1)")

    # ARIMA
    _card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(2.6),
          "ARIMA Forecasting", [
              "Stationarity: ADF + KPSS dual tests with auto-differencing",
              "Grid search: p ∈ {0..3}, d ∈ {0..2}, q ∈ {0..3} → 48 candidates",
              "Selection: AIC minimization (lower = better)",
              "Output: 30-day forecast + 95% confidence intervals",
              "Win rate: ~60% of stocks (vs ETS/linear)",
          ], ACCENT)

    # ETS
    _card(slide, Inches(0.5), Inches(4.3), Inches(5.8), Inches(2.2),
          "Holt-Winters Exponential Smoothing", [
              "ŷ(t+h) = lₜ + Σ φʲ · bₜ (damped additive trend)",
              "No seasonal component (daily stock data)",
              "Better for mean-reverting stocks (~30% win rate)",
              "Seasonal decomposition: STL with period=5",
          ], GREEN)

    # GARCH
    _card(slide, Inches(6.7), Inches(1.4), Inches(5.8), Inches(3.0),
          "GARCH(1,1) Volatility Model", [
              "σ²ₜ = ω + α·ε²ₜ₋₁ + β·σ²ₜ₋₁",
              "",
              "Captures volatility clustering (high vol → high vol)",
              "Persistence α + β ≈ 0.95 (slow mean reversion)",
              "Half-life: ~13.5 trading days",
              "30-day conditional volatility forecast",
              "MLE estimation via arch library (BFGS optimizer)",
          ], CYAN)

    # Regime table
    _card(slide, Inches(6.7), Inches(4.7), Inches(5.8), Inches(1.8),
          "Volatility Regime Classification", [
              "LOW (< 25th percentile)  →  Entry Signal",
              "NORMAL (25th–75th)  →  Neutral",
              "HIGH (75th–95th)  →  Caution",
              "EXTREME (≥ 95th)  →  Strong Caution",
          ], AMBER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — RISK & PORTFOLIO
# ═══════════════════════════════════════════════════════════════════════════════

def slide_risk_portfolio():
    slide = _add_slide()
    _section_header(slide, "Risk Analytics & Portfolio Optimization",
                    "Fama-French, Markowitz MVO, and Monte Carlo wealth planning")

    # Fama-French
    _card(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(2.4),
          "Fama-French 4-Factor Model", [
              "Rᵢ − Rf = α + β₁·MKT + β₂·SMB + β₃·HML + β₄·MOM + ε",
              "",
              "Market: Nifty 50 excess returns (Rf = 6.5%)",
              "SMB: Small − Big (size premium)",
              "HML: High − Low (value premium)",
              "MOM: Winners − Losers (momentum factor)",
          ], ACCENT)

    # Risk metrics
    _card(slide, Inches(0.5), Inches(4.1), Inches(6.0), Inches(1.5),
          "Risk Metrics Suite", [
              "Sharpe · Sortino · Beta · Jensen's Alpha · VaR(95%)",
              "Max Drawdown · Annualized Volatility",
              "Risk-free rate: 6.5% (10Y G-Sec yield)",
          ], GREEN)

    # Markowitz
    _card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.4),
          "Markowitz Mean-Variance Optimization", [
              "min w'Σw  s.t.  Σwᵢ = 1,  wᵢ ≤ 0.4",
              "",
              "2,000 Monte Carlo random portfolios",
              "SLSQP constrained optimization (scipy)",
              "Identifies: min-variance & max-Sharpe portfolios",
              "Efficient frontier scatter visualization",
          ], CYAN)

    # Monte Carlo
    _card(slide, Inches(6.8), Inches(4.1), Inches(5.8), Inches(2.5),
          "Monte Carlo Wealth Planner", [
              "1,000 simulations per scenario",
              "Monthly SIP with annual step-up support",
              "P10 / P25 / P50 / P75 / P90 wealth bands",
              "Goal-based reverse SIP solver (binary search)",
              "3 risk profiles: Conservative (9%) / Moderate (12%) / Aggressive (15%)",
          ], AMBER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — AI ADVISOR
# ═══════════════════════════════════════════════════════════════════════════════

def slide_ai_advisor():
    slide = _add_slide()
    _section_header(slide, "Agentic AI Advisor",
                    "The key differentiator — LLM with autonomous tool use")

    # Flow diagram
    flow_items = [
        ("User Query", '"Analyze Reliance"', ACCENT),
        ("Context Inject", "Profile + Portfolio\n+ Market Snapshot", GREEN),
        ("Llama 3.3 70B", "via Groq\n(500+ tok/s)", CYAN),
        ("SSE Stream", "Token-by-token\nto React frontend", ACCENT),
    ]
    for i, (title, sub, color) in enumerate(flow_items):
        lx = Inches(0.5 + i * 3.15)
        _add_rect(slide, lx, Inches(1.4), Inches(2.8), Inches(1.1),
                  WHITE, border_color=color, radius=0.06)
        _add_text_box(slide, lx + Inches(0.1), Inches(1.45), Inches(2.6), Inches(0.35),
                      title, font_size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
        _add_text_box(slide, lx + Inches(0.1), Inches(1.8), Inches(2.6), Inches(0.55),
                      sub, font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # Arrows between flow items
    for i in range(3):
        lx = Inches(3.4 + i * 3.15)
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      lx, Inches(1.75), Inches(0.35), Inches(0.25))
        arr.fill.solid()
        arr.fill.fore_color.rgb = GRAY
        arr.line.fill.background()

    # Tool box
    _add_rect(slide, Inches(3.5), Inches(2.8), Inches(6.3), Inches(1.5),
              LIGHT_AMBER, border_color=AMBER, radius=0.06)
    _add_text_box(slide, Inches(3.5), Inches(2.85), Inches(6.3), Inches(0.35),
                  "4 Live Function-Calling Tools (up to 3 autonomous rounds)",
                  font_size=12, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

    tools = [
        ("get_stock_analysis", "Technical signals, risk metrics, RF probability"),
        ("get_macro_dashboard", "VIX, USD/INR, Gold, Oil + market regime"),
        ("get_news_sentiment", "Headline sentiment scores for stocks"),
        ("run_screener", "Filter by signal, sector, Sharpe, drawdown"),
    ]
    for i, (name, desc) in enumerate(tools):
        col = i % 2
        row = i // 2
        lx = Inches(3.8 + col * 3.1)
        ty = Inches(3.3 + row * 0.45)
        _add_text_box(slide, lx, ty, Inches(1.5), Inches(0.35),
                      name, font_size=10, color=AMBER, bold=True)
        _add_text_box(slide, lx + Inches(1.5), ty, Inches(1.5), Inches(0.35),
                      desc, font_size=9, color=GRAY)

    # Bottom: what makes it agentic vs beyond RAG
    _card(slide, Inches(0.5), Inches(4.6), Inches(6.0), Inches(2.4),
          'What Makes It "Agentic"?', [
              "LLM autonomously decides when to call tools",
              "Multi-turn reasoning with tool results",
              "Portfolio-aware personalized advice",
              "Every recommendation backed by live quant data",
              "Inspired by Toolformer & ReAct frameworks",
          ], ACCENT)

    _card(slide, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.4),
          "Beyond Static RAG", [
              "Static RAG: inject data → generate (one-shot)",
              "Agentic: reason → decide tool → fetch → synthesize → respond",
              "",
              "Model: Llama 3.3 70B Versatile",
              "Provider: Groq LPU (500+ tokens/sec)",
              "First token: 2–5 seconds (SSE streaming)",
          ], CYAN)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — TECH STACK
# ═══════════════════════════════════════════════════════════════════════════════

def slide_tech_stack():
    slide = _add_slide()
    _section_header(slide, "Technology Stack & Frontend",
                    "Modern full-stack architecture with 18 interactive pages")

    # Backend
    _card(slide, Inches(0.5), Inches(1.3), Inches(3.8), Inches(3.3),
          "Backend", [
              "Python 3.12",
              "FastAPI 0.111",
              "pandas 2.2 + numpy 1.26",
              "scikit-learn 1.8",
              "statsmodels 0.14",
              "arch 8.0 (GARCH)",
              "scipy 1.17 (optimization)",
              "yfinance 0.2",
              "Pydantic v2 + httpx 0.27",
          ], ACCENT)

    # Frontend
    _card(slide, Inches(4.6), Inches(1.3), Inches(3.8), Inches(3.3),
          "Frontend", [
              "React 19",
              "TypeScript 5.9",
              "Vite 8.0 (build tool)",
              "Tailwind CSS 4",
              "Recharts 3.8 (7 chart types)",
              "Zustand 5.0 (state)",
              "TanStack React Query 5",
              "React Router 7",
              "Axios 1.13 + react-markdown",
          ], GREEN)

    # Infra
    _card(slide, Inches(8.7), Inches(1.3), Inches(3.8), Inches(3.3),
          "Infrastructure", [
              "Supabase (PostgreSQL + RLS)",
              "Hugging Face Spaces (Docker)",
              "Netlify CDN (frontend)",
              "Groq Cloud (LLM inference)",
              "APScheduler (daily cron)",
              "",
              "Docker: Python 3.12-slim",
              "Uvicorn ASGI server",
              "Port 7860 (HF Spaces)",
          ], CYAN)

    # 18 Pages — 3 columns
    _add_text_box(slide, Inches(0.5), Inches(4.9), Inches(4), Inches(0.35),
                  "18 Interactive Frontend Pages", font_size=16, color=PRIMARY, bold=True)

    page_groups = [
        ("Core", ACCENT, ["Dashboard", "AI Advisor Chat", "Portfolio Manager",
                          "Analytics Report", "Stock Screener", "Portfolio Optimizer"]),
        ("Advanced", CYAN, ["Time-Series Forecast", "Volatility Forecast", "Sector Rotation",
                            "Macro Dashboard", "Risk Factors", "ML Prediction"]),
        ("Planning", GREEN, ["Forward Planner", "Goal Planner", "Scenario Compare",
                             "News Sentiment", "Model Health (MLOps)", "Landing Page"]),
    ]

    for col, (group, color, pages) in enumerate(page_groups):
        lx = Inches(0.5 + col * 4.3)
        _add_text_box(slide, lx, Inches(5.3), Inches(2), Inches(0.3),
                      group, font_size=12, color=color, bold=True)
        tf = _add_rich_box(slide, lx, Inches(5.6), Inches(3.8), Inches(1.8))
        for page in pages:
            _add_para(tf, f"▸  {page}", size=10, color=DARK_GRAY, space_after=Pt(2))


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — RESULTS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_results():
    slide = _add_slide()
    _section_header(slide, "Results & Validation",
                    "Comprehensive testing across all system components")

    # KPI row
    kpis = [
        ("500+", "Stocks Analyzed\nDaily", ACCENT),
        ("35+", "API Endpoints\nTested & Passed", GREEN),
        ("< 100ms", "Cached Response\nLatency", CYAN),
        ("2–5s", "LLM First Token\n(Groq SSE)", AMBER),
    ]
    for i, (num, label, color) in enumerate(kpis):
        _kpi_card(slide, Inches(0.5 + i * 3.15), Inches(1.3),
                  Inches(2.9), Inches(1.3), num, label, color)

    # Model performance
    _card(slide, Inches(0.5), Inches(3.0), Inches(6.0), Inches(2.5),
          "Model Performance", [
              "RF Classifier accuracy: 52–58% (above 50% baseline)",
              "Regression ensemble dir. accuracy: 55–60%",
              "GARCH persistence (α+β): ~0.95",
              "ARIMA best-model rate: ~60% of stocks",
              "",
              "RSI vs TradingView: < 0.5% deviation",
              "MACD vs TradingView: < 0.3% deviation",
              "Prices vs Yahoo Finance: 0% deviation",
          ], ACCENT)

    # Testing results
    _card(slide, Inches(6.8), Inches(3.0), Inches(5.8), Inches(2.5),
          "Testing Summary", [
              "API endpoint tests: 35/35 passed ✓",
              "Edge case tests: 15/15 passed ✓",
              "Frontend page tests: 18/18 passed ✓",
              "Cross-browser: 6/6 browsers ✓",
              "Responsive breakpoints: 6/6 (320px → 2560px) ✓",
              "WCAG 2.1 AA accessibility: compliant ✓",
              "",
              "All 14 project objectives achieved ✓",
          ], GREEN)

    # Bottom bar
    _add_rect(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.55),
              PRIMARY, radius=0.08)
    _add_text_box(slide, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.45),
                  "Cold build: 30–60s  →  Cached: < 100ms  |  Pre-warm pipeline ensures zero cold-start for standard UX",
                  font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — COMPARISON
# ═══════════════════════════════════════════════════════════════════════════════

def slide_comparison():
    slide = _add_slide()
    _section_header(slide, "Comparison with Existing Platforms",
                    "PaisaPro.ai is the only platform combining all 12 capabilities")

    # Table
    headers = ["Capability", "PaisaPro", "Screener.in", "TradingView", "Groww", "ChatGPT"]
    rows = [
        ("Technical Signals (7 indicators)", "✓", "—", "✓", "—", "—"),
        ("ML Classification", "✓", "—", "—", "—", "—"),
        ("ML Regression Ensemble", "✓", "—", "—", "—", "—"),
        ("ARIMA / ETS Forecast", "✓", "—", "—", "—", "—"),
        ("GARCH Volatility", "✓", "—", "—", "—", "—"),
        ("Fama-French 4-Factor", "✓", "—", "—", "—", "—"),
        ("Portfolio Optimization", "✓", "—", "—", "—", "—"),
        ("Monte Carlo Planning", "✓", "—", "—", "—", "—"),
        ("AI Advisory (Live Data)", "✓", "—", "—", "—", "~"),
        ("Agentic Tool Use", "✓", "—", "—", "—", "—"),
        ("Portfolio Tracking", "✓", "—", "—", "✓", "—"),
        ("News Sentiment", "✓", "—", "—", "—", "—"),
    ]

    col_widths = [Inches(3.2), Inches(1.3), Inches(1.5), Inches(1.7), Inches(1.0), Inches(1.3)]
    table_left = Inches(1.2)
    table_top = Inches(1.4)
    row_height = Inches(0.42)

    # Header row
    hx = table_left
    for i, (hdr, cw) in enumerate(zip(headers, col_widths)):
        _add_rect(slide, hx, table_top, cw, row_height, ACCENT if i == 1 else PRIMARY, radius=0.02)
        _add_text_box(slide, hx, table_top + Inches(0.05), cw, Inches(0.3),
                      hdr, font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        hx += cw

    # Data rows
    for r, (cap, *vals) in enumerate(rows):
        ry = table_top + (r + 1) * row_height
        bg = LIGHT_BG if r % 2 == 0 else WHITE
        hx = table_left

        for i, (val, cw) in enumerate(zip([cap] + list(vals), col_widths)):
            _add_rect(slide, hx, ry, cw, row_height, bg, radius=0.0)
            if i == 0:
                _add_text_box(slide, hx + Inches(0.1), ry + Inches(0.05),
                              cw - Inches(0.2), Inches(0.3),
                              val, font_size=10, color=DARK_GRAY)
            else:
                color = GREEN if val == "✓" else (AMBER if val == "~" else RGBColor(200, 200, 200))
                _add_text_box(slide, hx, ry + Inches(0.02), cw, Inches(0.35),
                              val, font_size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
            hx += cw

    # Callout
    _add_rect(slide, Inches(2.5), Inches(6.7), Inches(8.3), Inches(0.5),
              LIGHT_ACCENT, border_color=ACCENT, radius=0.1)
    _add_text_box(slide, Inches(2.5), Inches(6.73), Inches(8.3), Inches(0.4),
                  "PaisaPro.ai: 12/12 capabilities  |  Closest competitor: 1–2/12",
                  font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — LIMITATIONS & FUTURE WORK
# ═══════════════════════════════════════════════════════════════════════════════

def slide_limitations():
    slide = _add_slide()
    _section_header(slide, "Limitations & Future Work")

    # Limitations
    _card(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(3.2),
          "Current Limitations", [
              "1.  Single-user mode (hardcoded user ID)",
              "2.  Daily OHLCV only (no intraday data)",
              "3.  Keyword-based sentiment (not FinBERT)",
              "4.  Proxy Fama-French factors (not academic factor portfolios)",
              "5.  Normal distribution assumption in Monte Carlo",
              "6.  LLM hallucination risk (despite context grounding)",
              "7.  No comprehensive walk-forward portfolio backtest",
          ], RED)

    # Future work
    _card(slide, Inches(6.7), Inches(1.3), Inches(5.8), Inches(5.3),
          "Proposed Future Enhancements", [
              "1.  JWT multi-user authentication (Supabase Auth)",
              "2.  FinBERT transformer-based sentiment analysis",
              "3.  Intraday data support (5-min / 15-min candles)",
              "4.  Fama-French 5-factor model (+ Quality, Investment)",
              "5.  Options analytics (Greeks, implied volatility surface)",
              "6.  Walk-forward portfolio backtesting with tx costs",
              "7.  React Native mobile application",
              "8.  Real-time WebSocket price feeds",
              "9.  Fundamental analysis (P/E, P/B, debt ratios)",
              "10. SEBI RIA compliance exploration",
          ], GREEN)

    # Lessons learned
    _card(slide, Inches(0.5), Inches(4.8), Inches(5.8), Inches(2.0),
          "Key Lessons Learned", [
              "Caching is critical — 30s → <100ms with TTL cache",
              "Model simplicity wins for daily data (250–1250 obs)",
              "Agentic prompt engineering requires iteration",
              "SSE streaming transforms perceived performance",
          ], AMBER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════════

def slide_conclusion():
    slide = _add_slide()
    _section_header(slide, "Conclusion")

    # Big KPIs
    kpis = [
        ("500+", "NSE Stocks\nAnalyzed Daily", ACCENT),
        ("19", "Backend Service\nModules", GREEN),
        ("4", "Agentic AI\nLive Tools", CYAN),
        ("14/14", "Project Objectives\nAchieved", AMBER),
    ]
    for i, (num, label, color) in enumerate(kpis):
        _kpi_card(slide, Inches(0.5 + i * 3.15), Inches(1.3),
                  Inches(2.9), Inches(1.5), num, label, color)

    # Key contributions
    _card(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(2.5),
          "Key Contributions", [
              "1.  First platform to integrate ARIMA, GARCH, Fama-French, Markowitz, Monte Carlo, ML classification, ML regression, and agentic LLM advisory in one system for Indian retail investors",
              "2.  Novel agentic AI advisor with 4 live function-calling tools — moves beyond static RAG toward autonomous multi-step reasoning",
              "3.  Robust 3-model ML regression ensemble with OHLC microstructure features, vol-adjusted targets, and walk-forward validation",
              "4.  Production-grade architecture: 19 independent modules, 5-tier TTL cache, daily auto-refresh pipeline, Docker deployment",
              "5.  18-page responsive React frontend making institutional-grade analytics accessible to retail investors",
          ], ACCENT, body_size=11)

    # Bottom statement
    _add_rect(slide, Inches(1.5), Inches(6.0), Inches(10.3), Inches(1.0),
              PRIMARY, radius=0.1)
    _add_text_box(slide, Inches(1.5), Inches(6.1), Inches(10.3), Inches(0.4),
                  "PaisaPro.ai demonstrates that institutional-grade financial analytics",
                  font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.4),
                  "can be democratized for 150M+ Indian retail investors",
                  font_size=16, color=CYAN, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 17 — THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════

def slide_thankyou():
    slide = _add_slide()
    _set_bg(slide, PRIMARY)

    # Decorative circles
    for cx, cy, sz in [(0.3, 5.5, 3.5), (10.5, 0.3, 3), (12, 5, 2)]:
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(cx), Inches(cy),
                                    Inches(sz), Inches(sz))
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT
        c.fill.fore_color.brightness = 0.7
        c.line.fill.background()

    # Logo
    logo_path = os.path.join(os.path.dirname(__file__), "NMIMS logo.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(5.9), Inches(0.8),
                                  height=Inches(0.8))

    _add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(0.9),
                  "Thank You!", font_size=52, color=WHITE, bold=True,
                  align=PP_ALIGN.CENTER)

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(4.7), Inches(3.0),
                                   Inches(4), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()

    _add_text_box(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(0.6),
                  "Questions & Discussion", font_size=26, color=WHITE,
                  align=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.4),
                  "Stephen Baraik  ·  Tanisha Ghosh  ·  Sneha Das",
                  font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.3),
                  "Supervisor: Dr. Suresh Pathare  |  Associate Professor, Data Science",
                  font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(0.3),
                  "SVKM's NMIMS University  ·  SOMASA  ·  2025–2026",
                  font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  BUILD
# ═══════════════════════════════════════════════════════════════════════════════

TOTAL_SLIDES = 17

slide_title()        # 1
slide_agenda()       # 2
slide_problem()      # 3
slide_solution()     # 4
slide_architecture() # 5
slide_data_pipeline()# 6
slide_analytics()    # 7
slide_ml()           # 8
slide_forecasting()  # 9
slide_risk_portfolio()# 10
slide_ai_advisor()   # 11
slide_tech_stack()   # 12
slide_results()      # 13
slide_comparison()   # 14
slide_limitations()  # 15
slide_conclusion()   # 16
slide_thankyou()     # 17

# Add slide numbers (skip title and thank-you)
for i, slide in enumerate(prs.slides):
    if i > 0 and i < len(prs.slides) - 1:
        _add_slide_number(slide, i + 1, TOTAL_SLIDES)

out_path = os.path.join(os.path.dirname(__file__), "PaisaPro_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
